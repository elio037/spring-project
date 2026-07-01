# -*- coding: utf-8 -*-
"""BookReview 기술 리포트 PPT 생성"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ── 팔레트 ──
INK    = RGBColor(0x16, 0x1B, 0x2E)   # deep navy bg
INK2   = RGBColor(0x20, 0x27, 0x42)   # card on dark
BLUE   = RGBColor(0x2E, 0x7C, 0xE6)   # primary blue
GOLD   = RGBColor(0xF5, 0xA6, 0x23)   # star gold accent
LIGHT  = RGBColor(0xF4, 0xF5, 0xF8)   # light content bg
CARD   = RGBColor(0xFF, 0xFF, 0xFF)   # white card
INKTXT = RGBColor(0x1D, 0x1D, 0x1F)   # dark text
GREY   = RGBColor(0x6A, 0x70, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTE   = RGBColor(0xB6, 0xBE, 0xD0)   # muted on dark
GREEN  = RGBColor(0x2E, 0xA0, 0x6A)
RED    = RGBColor(0xE0, 0x4A, 0x4A)

KFONT = '맑은 고딕'

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide(bg):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    # 배경을 맨 뒤로
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s

def _set_font(run, size, color, bold, font, italic=False):
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', font)

def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    """runs: list of (text,size,color,bold[,italic]) or list-of-list for paragraphs"""
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for seg in para:
            txt, size, color, bold = seg[0], seg[1], seg[2], seg[3]
            italic = seg[4] if len(seg) > 4 else False
            r = p.add_run(); r.text = txt
            _set_font(r, size, color, bold, KFONT, italic)
    return tb

def card(s, l, t, w, h, fill, radius=0.10, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    try: shp.adjustments[0] = radius
    except Exception: pass
    return shp

def circle(s, l, t, d, fill, label=None, lsize=18, lcolor=WHITE):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    c.fill.solid(); c.fill.fore_color.rgb = fill; c.line.fill.background()
    c.shadow.inherit = False
    if label is not None:
        tf = c.text_frame; tf.word_wrap = False
        tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label; _set_font(r, lsize, lcolor, True, KFONT)
    return c

def chip(s, l, t, w, h, txt, fill, tcolor):
    c = card(s, l, t, w, h, fill, radius=0.5)
    tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt; _set_font(r, 11.5, tcolor, True, KFONT)
    return c

def kicker(s, l, t, txt, color=GOLD):
    text(s, l, t, Inches(6), Inches(0.35), [(txt, 13, color, True)])

# ════════════════════ 1. 표지 (dark) ════════════════════
s = slide(INK)
circle(s, Inches(0.9), Inches(0.85), Inches(0.62), GOLD, '★', 24, INK)
text(s, Inches(1.65), Inches(0.92), Inches(6), Inches(0.5), [('BookReview', 18, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.4),
     [('스프링 발표자료', 54, WHITE, True)])
text(s, Inches(0.95), Inches(3.75), Inches(11), Inches(0.7),
     [('독자 리뷰 기반 도서 탐색 플랫폼 — Spring MVC 발표자료', 20, MUTE, False)])
# 하단 메타 칩
chip(s, Inches(0.95), Inches(5.6), Inches(2.7), Inches(0.5), 'Spring MVC · MyBatis', INK2, MUTE)
chip(s, Inches(3.8),  Inches(5.6), Inches(1.8), Inches(0.5), 'Oracle DB', INK2, MUTE)
chip(s, Inches(5.75), Inches(5.6), Inches(2.9), Inches(0.5), 'face-api.js 얼굴 로그인', INK2, GOLD)
text(s, Inches(0.95), Inches(6.55), Inches(11), Inches(0.4),
     [('kr.ac.kopo  ·  2026  ·  v1.0', 12, GREY, False)])

# ════════════════════ 2. 프로젝트 개요 ════════════════════
s = slide(LIGHT)
kicker(s, Inches(0.7), Inches(0.5), '01  OVERVIEW')
text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8), [('프로젝트 개요', 38, INKTXT, True)])
# 한 줄 정의 배너 — "무엇을 하는 사이트인가"를 즉시 인지시키는 핵심 문장
card(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(0.95), BLUE, radius=0.14)
text(s, Inches(1.1), Inches(1.8), Inches(11.2), Inches(0.95),
     [[('“독자가 남긴 실제 리뷰”로 책을 고르는  ', 19, WHITE, True),
       ('도서 리뷰 · 탐색 플랫폼', 19, RGBColor(0xFF,0xE6,0xB0), True)]],
     anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.7), Inches(3.0), Inches(12), Inches(1.0),
     [[('알라딘 도서 리뷰 크롤링 데이터', 15, BLUE, True), ('를 정제해 Oracle에 적재하고, ', 15, INKTXT, False),
       ('도서 검색 · 평점/리뷰 순위 · 상세 통계(평점 분포·감성 분석)', 15, INKTXT, False)],
      [('를 제공하는 Spring MVC 웹 애플리케이션. ', 15, INKTXT, False), ('얼굴 인식 로그인', 15, INKTXT, True),
       ('과 카카오 로그인까지 지원한다.', 15, INKTXT, False)]],
     line_spacing=1.25, space_after=6)
# 스탯 콜아웃 카드 4개
stats = [('200', '도서', BLUE), ('4,000', '샘플 리뷰', GOLD), ('22', 'MVC 엔드포인트', GREEN), ('128', '얼굴 특징 차원', RED)]
x = Inches(0.7)
for val, lab, col in stats:
    card(s, x, Inches(4.05), Inches(2.85), Inches(2.4), CARD, radius=0.10)
    text(s, x, Inches(4.45), Inches(2.85), Inches(1.0), [(val, 46, col, True)], align=PP_ALIGN.CENTER)
    text(s, x, Inches(5.55), Inches(2.85), Inches(0.5), [(lab, 14, GREY, True)], align=PP_ALIGN.CENTER)
    x = Emu(int(x) + int(Inches(3.05)))

# ════════════════════ 3. 기술 스택 ════════════════════
s = slide(LIGHT)
kicker(s, Inches(0.7), Inches(0.5), '02  TECH STACK')
text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8), [('기술 스택', 38, INKTXT, True)])
groups = [
    ('Backend', BLUE, ['Java 17', 'Spring Web MVC 6.2', 'MyBatis 3.5 + mybatis-spring']),
    ('Database', GOLD, ['Oracle XE 21c', 'ojdbc11 / commons-dbcp', 'BR_BOOK·REVIEW·MEMBER', '+ COMMENT·LIKE·WISHLIST']),
    ('Frontend', GREEN, ['JSP / JSTL (Jakarta)', '자체 CSS 디자인 시스템', '반응형 카드 레이아웃']),
    ('AI / 인증', RED, ['딥러닝 얼굴 인식 (face-api.js)', '모델 3종 — 검출·특징 추출', '유클리드 거리로 얼굴 비교']),
]
positions = [(0.7, 1.9), (6.95, 1.9), (0.7, 4.35), (6.95, 4.35)]
for (gx, gy), (name, col, items) in zip(positions, groups):
    card(s, Inches(gx), Inches(gy), Inches(5.65), Inches(2.2), CARD, radius=0.08)
    circle(s, Inches(gx+0.3), Inches(gy+0.32), Inches(0.45), col)
    text(s, Inches(gx+0.95), Inches(gy+0.3), Inches(4.5), Inches(0.5),
         [(name, 19, INKTXT, True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(gx+0.95), Inches(gy+1.0), Inches(4.5), Inches(1.1),
         [[(('•  '+it), 13, GREY, False)] for it in items], line_spacing=1.15, space_after=3)

# ════════════════════ 4. 시스템 아키텍처 ════════════════════
s = slide(INK)
kicker(s, Inches(0.7), Inches(0.5), '03  ARCHITECTURE')
text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8), [('시스템 아키텍처', 38, WHITE, True)])
text(s, Inches(0.7), Inches(1.8), Inches(11.8), Inches(0.5),
     [('표준 레이어드 아키텍처 — 요청은 DispatcherServlet을 통해 계층을 따라 흐른다.', 15, MUTE, False)])
# (name, 한글 역할=발표용 설명, 기술 태그, 색)
layers = [('브라우저',             '화면 표시 · 입력',  'JSP',         BLUE),
          ('Dispatcher\nServlet',  '요청 배분(관문)',   'spring-mvc',  GOLD),
          ('Controller',           '요청 받고 응답',    '@Controller', GREEN),
          ('Service',              '핵심 로직 처리',    '@Service',    BLUE),
          ('DAO + Mapper',         'SQL 실행',         'MyBatis',     GOLD),
          ('Oracle',               '데이터 저장',       'BR_* 테이블',  RED)]
n = len(layers); bx = Inches(0.7); bw = Inches(1.85); gap = Inches(0.16); by = Inches(2.75); bh = Inches(1.9)
x = bx
for i, (name, role, tech, col) in enumerate(layers):
    c = card(s, x, by, bw, bh, INK2, radius=0.12)
    circle(s, Emu(int(x)+int(Inches(0.72))), Emu(int(by)+int(Inches(0.2))), Inches(0.4), col)
    text(s, x, Emu(int(by)+int(Inches(0.64))), bw, Inches(0.42),
         [(name.replace('\n', ' '), 14, WHITE, True)], align=PP_ALIGN.CENTER, line_spacing=0.95)
    text(s, x, Emu(int(by)+int(Inches(1.06))), bw, Inches(0.38),
         [(role, 12, RGBColor(0xE6, 0xEA, 0xF5), True)], align=PP_ALIGN.CENTER)
    text(s, x, Emu(int(by)+int(Inches(1.46))), bw, Inches(0.3),
         [(tech, 9, MUTE, False)], align=PP_ALIGN.CENTER)
    if i < n-1:
        text(s, Emu(int(x)+int(bw)-int(Inches(0.02))), Emu(int(by)+int(Inches(0.75))),
             Inches(0.3), Inches(0.5), [('▸', 20, GOLD, True)], align=PP_ALIGN.CENTER)
    x = Emu(int(x) + int(bw) + int(gap))
# 요청 흐름 한 줄 설명(발표용 스크립트)
text(s, Inches(0.7), Inches(4.95), Inches(12.4), Inches(0.5),
     [[('요청 흐름:  ', 13, GOLD, True),
       ('브라우저 요청 → DispatcherServlet이 알맞은 Controller로 전달 → Service 로직 → DAO·Mapper가 SQL 실행 → Oracle 응답',
        13, MUTE, False)]])
text(s, Inches(0.7), Inches(5.55), Inches(12), Inches(0.6),
     [[('패키지: ', 14, GOLD, True), ('kr.ac.kopo.{main · book · member · member.social · comment · wishlist} × (controller·service·dao·vo)', 14, MUTE, False)]])
text(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.6),
     [[('인코딩: ', 14, GOLD, True), ('CharacterEncodingFilter(UTF-8) 전역 적용', 14, MUTE, False)]])

# ════════════════════ 5. DB 설계 ════════════════════
s = slide(LIGHT)
kicker(s, Inches(0.7), Inches(0.5), '04  DATABASE')
text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8), [('데이터베이스 설계', 38, INKTXT, True)])
tables = [
    ('BR_MEMBER', BLUE, '회원', ['NO (PK)', 'ID / PASSWORD', 'NICKNAME / ROLE', 'FACE_DESCRIPTOR', 'PROVIDER / PROVIDER_ID']),
    ('BR_BOOK', GOLD, '도서 집계', ['NO (PK)', 'TITLE', 'REVIEW_COUNT', 'AVG_RATING', 'DIST1~5', 'POS/NEU/NEG_COUNT']),
    ('BR_REVIEW', GREEN, '샘플 리뷰', ['NO (PK)', 'BOOK_NO (FK)', 'CONTENT', 'RATING']),
]
gx = 0.7
for name, col, desc, cols in tables:
    card(s, Inches(gx), Inches(1.95), Inches(3.95), Inches(4.4), CARD, radius=0.06)
    head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(gx), Inches(1.95), Inches(3.95), Inches(0.85))
    head.fill.solid(); head.fill.fore_color.rgb = col; head.line.fill.background(); head.shadow.inherit=False
    head.adjustments[0] = 0.18
    text(s, Inches(gx+0.3), Inches(2.02), Inches(3.4), Inches(0.45), [(name, 17, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(gx+0.3), Inches(2.4), Inches(3.4), Inches(0.35), [(desc, 11, RGBColor(0xFF,0xFF,0xFF), False)])
    text(s, Inches(gx+0.35), Inches(3.1), Inches(3.3), Inches(3.0),
         [[(('•  '+c), 13, INKTXT, ('PK' in c or 'FK' in c or 'FACE' in c))] for c in cols],
         line_spacing=1.2, space_after=5)
    gx += 4.15
text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
     [[('FACE_DESCRIPTOR', 13, RED, True), (' = 128차원 얼굴 특징 벡터를 JSON 문자열로 저장  ·  FK: BR_REVIEW.BOOK_NO → BR_BOOK.NO', 13, GREY, False)]])

# ════════════════════ 5B. 확장 데이터 모델 ════════════════════
s = slide(LIGHT)
kicker(s, Inches(0.7), Inches(0.5), 'NEW  ·  EXTENDED SCHEMA')
text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.8), [('확장 데이터 모델 — 리뷰 · 좋아요 · 찜', 33, INKTXT, True)])
text(s, Inches(0.7), Inches(1.72), Inches(12), Inches(0.5),
     [('회원 활동 기능을 위해 테이블 3개를 추가했다. 회원 연관 데이터는 탈퇴 시 함께 삭제된다.', 14, GREY, False)])
ntables = [
    ('BR_COMMENT', BLUE, '회원 리뷰', ['NO (PK)', 'MEMBER_NO (FK)', 'BOOK_NO (FK)', 'CONTENT', 'RATING (1~5)', 'REG_DATE']),
    ('BR_COMMENT_LIKE', GREEN, '리뷰 좋아요', ['NO (PK)', 'MEMBER_NO (FK)', 'COMMENT_NO (FK)', 'UNIQUE (회원, 리뷰)']),
    ('BR_WISHLIST', GOLD, '찜', ['NO (PK)', 'MEMBER_NO (FK)', 'BOOK_NO (FK)', 'UNIQUE (회원, 도서)']),
]
gx = 0.7
for name, col, desc, cols in ntables:
    card(s, Inches(gx), Inches(2.35), Inches(3.95), Inches(3.75), CARD, radius=0.06)
    head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(gx), Inches(2.35), Inches(3.95), Inches(0.85))
    head.fill.solid(); head.fill.fore_color.rgb = col; head.line.fill.background(); head.shadow.inherit = False
    head.adjustments[0] = 0.18
    text(s, Inches(gx+0.3), Inches(2.42), Inches(3.4), Inches(0.45), [(name, 15, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(gx+0.3), Inches(2.8), Inches(3.4), Inches(0.35), [(desc, 11, WHITE, False)])
    text(s, Inches(gx+0.35), Inches(3.5), Inches(3.3), Inches(2.4),
         [[(('•  '+c), 12.5, INKTXT, ('PK' in c or 'FK' in c or 'UNIQUE' in c))] for c in cols],
         line_spacing=1.2, space_after=5)
    gx += 4.15
text(s, Inches(0.7), Inches(6.45), Inches(12.2), Inches(0.5),
     [[('ON DELETE CASCADE', 13, RED, True),
       ('  —  회원 탈퇴 시 작성 리뷰·좋아요·찜 자동 삭제   ·   UNIQUE 제약으로 중복 좋아요/찜 방지', 13, GREY, False)]])

# ════════════════════ 5C. ERD (전체 관계도) ════════════════════
s = slide(LIGHT)
kicker(s, Inches(0.7), Inches(0.42), 'ERD  ·  ENTITY RELATIONSHIP')
text(s, Inches(0.7), Inches(0.74), Inches(12), Inches(0.7), [('ERD — 전체 테이블 관계도', 31, INKTXT, True)])
# matplotlib로 렌더한 ERD 이미지(설명 패널 포함) 삽입 — 비율 0.616
import os
_erd = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'erd_bookreview.png')
s.shapes.add_picture(_erd, Inches(1.92), Inches(1.45), width=Inches(9.5))

# ════════════════════ 6. 주요 기능 ════════════════════
s = slide(LIGHT)
kicker(s, Inches(0.7), Inches(0.5), '05  FEATURES')
text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8), [('주요 기능', 38, INKTXT, True)])
feats = [
    ('🔍', '도서 검색', '제목 일부로 검색, 리뷰 많은 순 최대 50건', BLUE),
    ('🏆', '평점 순위', '리뷰 N개 이상 기준 평점 평균 TOP 100', GOLD),
    ('📊', '리뷰 순위', '리뷰 수 기준 인기 도서 TOP 100', GREEN),
    ('📖', '도서 상세', '평점 분포 · 감성 비율 · 샘플 리뷰', BLUE),
    ('👤', '회원/로그인', '가입 검증 · 세션 인증 · 중복 검사', RED),
    ('😊', '얼굴 로그인', '카메라로 즉시 로그인 (거리 매칭)', GOLD),
]
i = 0
for row in range(2):
    for coli in range(3):
        fx = 0.7 + coli*4.15; fy = 1.95 + row*2.25
        emoji, name, desc, col = feats[i]; i += 1
        card(s, Inches(fx), Inches(fy), Inches(3.95), Inches(2.0), CARD, radius=0.08)
        circle(s, Inches(fx+0.32), Inches(fy+0.32), Inches(0.7), col, emoji, 22, WHITE)
        text(s, Inches(fx+1.25), Inches(fy+0.45), Inches(2.6), Inches(0.6), [(name, 18, INKTXT, True)], anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(fx+0.35), Inches(fy+1.25), Inches(3.3), Inches(0.7), [(desc, 12.5, GREY, False)], line_spacing=1.15)

# ════════════════════ 6B. 회원 기능 확장 ════════════════════
s = slide(LIGHT)
kicker(s, Inches(0.7), Inches(0.5), 'NEW  ·  MEMBER FEATURES')
text(s, Inches(0.7), Inches(0.85), Inches(11.5), Inches(0.8), [('회원 기능 확장', 38, INKTXT, True)])
nfeats = [
    ('💬', '카카오 로그인', 'OAuth2 인가코드\n최초 시 자동 가입', GOLD),
    ('⭐', '별점 리뷰', '1~5점 + 본문\n독자 리뷰에 통합', BLUE),
    ('❤', '리뷰 좋아요', '토글 · 좋아요 수\n도움된 리뷰 표시', RED),
    ('🔖', '찜하기', '보고싶은 책 저장\n마이페이지 내 서재', GREEN),
    ('🙋', '마이페이지', '내 리뷰·좋아요·찜\n회원 탈퇴', BLUE),
]
x = 0.7
for emoji, name, desc, col in nfeats:
    card(s, Inches(x), Inches(2.2), Inches(2.35), Inches(3.45), CARD, radius=0.09)
    circle(s, Inches(x+0.82), Inches(2.55), Inches(0.7), col, emoji, 20, WHITE)
    text(s, Inches(x), Inches(3.45), Inches(2.35), Inches(0.5), [(name, 16, INKTXT, True)], align=PP_ALIGN.CENTER)
    text(s, Inches(x+0.2), Inches(4.15), Inches(1.95), Inches(1.3), [(desc, 12, GREY, False)], align=PP_ALIGN.CENTER, line_spacing=1.2)
    x += 2.5
text(s, Inches(0.7), Inches(6.15), Inches(12.2), Inches(0.5),
     [[('Spring Security 없이 ', 13, GREY, False), ('세션 기반 인증', 13, BLUE, True),
       ('으로 구현   ·   JDK HttpClient로 카카오 OAuth 직접 연동', 13, GREY, False)]])

# ════════════════════ 6C. 권한별 기능 (비회원/회원/관리자) ════════════════════
s = slide(LIGHT)
kicker(s, Inches(0.7), Inches(0.5), 'ROLES  ·  PERMISSIONS')
text(s, Inches(0.7), Inches(0.85), Inches(12.2), Inches(0.8),
     [('사용자 권한별 기능 — 비회원 · 회원 · 관리자', 30, INKTXT, True)])
text(s, Inches(0.7), Inches(1.65), Inches(12.2), Inches(0.4),
     [[('세션 기반 인증 · ', 13, GREY, False), ('ROLE 컬럼', 13, BLUE, True),
       ('으로 일반 회원과 관리자(ADMIN)를 구분한다.', 13, GREY, False)]])
roles = [
    ('🌐  비회원 (Guest)', GREY, '로그인 없이 열람만', [
        '도서 검색 (제목 일부)',
        '평점 순위 · 리뷰 순위 보기',
        '도서 상세 — 평점 분포·감성·샘플 리뷰',
        '— 리뷰 작성·찜·좋아요 불가',
    ]),
    ('👤  회원 (User)', BLUE, '비회원 기능 + 활동', [
        '회원가입 · 로그인 / 카카오 로그인',
        '얼굴 등록 · 얼굴 로그인',
        '별점 리뷰 작성 · 삭제(본인)',
        '리뷰 좋아요(토글) · 찜하기',
        '마이페이지 — 내 리뷰·좋아요·찜·탈퇴',
    ]),
    ('🛡  관리자 (Admin)', RED, '회원 기능 + 운영', [
        '관리자 대시보드 — 회원 수·도서 수',
        '전체 회원 목록 조회',
        '도서 표지 이미지 경로 관리·수정',
        'ROLE=ADMIN 세션만 /admin 접근',
    ]),
]
x = 0.7
for title_, col, sub, items in roles:
    card(s, Inches(x), Inches(2.25), Inches(3.95), Inches(4.55), CARD, radius=0.06)
    head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.25), Inches(3.95), Inches(0.95))
    head.fill.solid(); head.fill.fore_color.rgb = col; head.line.fill.background(); head.shadow.inherit = False
    head.adjustments[0] = 0.12
    text(s, Inches(x+0.3), Inches(2.38), Inches(3.4), Inches(0.45), [(title_, 16, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(x+0.3), Inches(2.82), Inches(3.4), Inches(0.35), [(sub, 11, WHITE, False)])
    text(s, Inches(x+0.32), Inches(3.45), Inches(3.45), Inches(3.2),
         [[((('•  '+it) if not it.startswith('—') else it), 12,
            (RED if it.startswith('—') else GREY), False)] for it in items],
         line_spacing=1.2, space_after=7)
    x += 4.15

# ════════════════════ 6D. 로그인 기능 ════════════════════
s = slide(LIGHT)
kicker(s, Inches(0.7), Inches(0.5), 'LOGIN  ·  AUTHENTICATION')
text(s, Inches(0.7), Inches(0.85), Inches(12.2), Inches(0.8), [('로그인 — 3가지 방식', 33, INKTXT, True)])
text(s, Inches(0.7), Inches(1.65), Inches(12.2), Inches(0.4),
     [[('세션 기반 인증.  ', 13, GREY, False), ('카카오 소셜 로그인', 13, GOLD, True), ('과 ', 13, GREY, False),
       ('얼굴 인식 로그인', 13, BLUE, True), ('까지 지원한다.', 13, GREY, False)]])
logins = [
    ('🔑  아이디 · 비밀번호', GREEN, '기본 로그인', [
        '회원가입 시 입력값 검증',
        '세션(HttpSession) 인증',
        '로그아웃 · 회원 탈퇴',
    ]),
    ('💬  카카오 로그인', GOLD, 'OAuth 2.0', [
        '카카오 인가코드 방식',
        '최초 로그인 시 자동 회원가입',
        '비밀번호 없이 소셜 인증',
    ]),
    ('😊  얼굴 인식 로그인', BLUE, '딥러닝 (face-api.js)', [
        '카메라로 얼굴 등록·로그인',
        '128차원 얼굴 특징 추출',
        '유클리드 거리로 본인 확인',
    ]),
]
x = 0.7
for title_, col, tag, items in logins:
    card(s, Inches(x), Inches(2.25), Inches(3.95), Inches(4.3), CARD, radius=0.06)
    head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.25), Inches(3.95), Inches(0.95))
    head.fill.solid(); head.fill.fore_color.rgb = col; head.line.fill.background(); head.shadow.inherit = False
    head.adjustments[0] = 0.12
    text(s, Inches(x+0.3), Inches(2.4), Inches(3.4), Inches(0.45), [(title_, 15.5, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(x+0.3), Inches(2.82), Inches(3.4), Inches(0.35), [(tag, 11, WHITE, False)])
    text(s, Inches(x+0.35), Inches(3.5), Inches(3.4), Inches(2.8),
         [[(('•  '+it), 13, INKTXT, False)] for it in items], line_spacing=1.35, space_after=9)
    x += 4.15

# ════════════════════ 7. 감성 분석 ════════════════════
s = slide(INK)
kicker(s, Inches(0.7), Inches(0.5), '06  SENTIMENT')
text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8), [('평점 분포 & 감성 분석', 38, WHITE, True)])
text(s, Inches(0.7), Inches(1.85), Inches(6.0), Inches(2.0),
     [[('리뷰 평점(1~5)을 기준으로 감성을 ', 15, MUTE, False)],
      [('단순·명확하게 3분류', 15, GOLD, True), ('하여', 15, MUTE, False)],
      [('도서 상세에서 막대로 시각화한다.', 15, MUTE, False)]],
     line_spacing=1.3, space_after=6)
rules = [('긍정', '4 ~ 5점', GREEN), ('중립', '3점', GOLD), ('부정', '1 ~ 2점', RED)]
ry = 3.9
for name, rng, col in rules:
    circle(s, Inches(0.7), Inches(ry), Inches(0.5), col)
    text(s, Inches(1.4), Inches(ry+0.02), Inches(2), Inches(0.5), [(name, 18, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.2), Inches(ry+0.02), Inches(2.5), Inches(0.5), [(rng, 16, MUTE, False)], anchor=MSO_ANCHOR.MIDDLE)
    ry += 0.85
# 오른쪽: 평점 분포 막대 시각화 (예시 - 데미안)
card(s, Inches(7.1), Inches(1.95), Inches(5.5), Inches(4.7), INK2, radius=0.06)
text(s, Inches(7.45), Inches(2.2), Inches(5), Inches(0.5), [[('예시 — ', 14, MUTE, False), ('데미안', 14, GOLD, True), ('  (평균 4.63 ★)', 14, MUTE, False)]])
dist = [('5★', 419, GREEN), ('4★', 102, GREEN), ('3★', 20, GOLD), ('2★', 12, RED), ('1★', 8, RED)]
maxv = 419; barx = 8.4; topy = 3.0; bw_full = 2.9
for lab, v, col in dist:
    text(s, Inches(7.45), Inches(topy), Inches(0.9), Inches(0.4), [(lab, 13, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    track = card(s, Inches(barx), Inches(topy+0.05), Inches(bw_full), Inches(0.34), RGBColor(0x33,0x3B,0x55), radius=0.5)
    fillw = max(0.18, bw_full * v / maxv)
    fill = card(s, Inches(barx), Inches(topy+0.05), Inches(fillw), Inches(0.34), col, radius=0.5)
    text(s, Inches(barx+bw_full+0.1), Inches(topy), Inches(0.9), Inches(0.4), [(str(v), 12, MUTE, False)], anchor=MSO_ANCHOR.MIDDLE)
    topy += 0.62
text(s, Inches(7.45), Inches(6.05), Inches(5), Inches(0.5),
     [[('긍정 ', 13, GREEN, True), ('92%', 13, GREEN, True), ('  ·  중립 ', 13, GOLD, True), ('4%', 13, GOLD, True), ('  ·  부정 ', 13, RED, True), ('4%', 13, RED, True)]])

# ════════════════════ 8. 얼굴 로그인 (핵심) ════════════════════
s = slide(INK)
kicker(s, Inches(0.7), Inches(0.5), '07  FACE LOGIN')
text(s, Inches(0.7), Inches(0.85), Inches(11.5), Inches(0.8), [('핵심 기술 — 얼굴 인식 로그인', 36, WHITE, True)])
text(s, Inches(0.7), Inches(1.8), Inches(12), Inches(0.5),
     [('얼굴 인식은 브라우저(face-api.js)에서 수행하고, 서버는 특징 벡터 저장·조회만 담당한다.', 14, MUTE, False)])
steps = [
    ('1', '얼굴 감지', 'face-api.js가 카메라\n프레임에서 얼굴 검출'),
    ('2', '특징 추출', '128차원 descriptor\n벡터 생성'),
    ('3', '거리 비교', '등록된 전 회원과\n유클리드 거리 계산'),
    ('4', '인증', '최소 거리 ≤ 0.5 →\n해당 회원 자동 로그인'),
]
x = 0.7; cols_ = [BLUE, GOLD, GREEN, RED]
for i, (num, name, desc) in enumerate(steps):
    card(s, Inches(x), Inches(2.7), Inches(2.85), Inches(2.7), INK2, radius=0.10)
    circle(s, Inches(x+1.07), Inches(2.95), Inches(0.7), cols_[i], num, 24, WHITE)
    text(s, Inches(x), Inches(3.85), Inches(2.85), Inches(0.5), [(name, 18, WHITE, True)], align=PP_ALIGN.CENTER)
    text(s, Inches(x+0.2), Inches(4.4), Inches(2.45), Inches(0.9), [(desc, 12.5, MUTE, False)], align=PP_ALIGN.CENTER, line_spacing=1.1)
    if i < 3:
        text(s, Inches(x+2.78), Inches(3.75), Inches(0.35), Inches(0.5), [('▸', 22, GOLD, True)], align=PP_ALIGN.CENTER)
    x += 3.05
text(s, Inches(0.7), Inches(5.85), Inches(12), Inches(0.7),
     [[('모델: ', 14, GOLD, True), ('ssdMobilenetv1 + faceLandmark68Net + faceRecognitionNet', 14, MUTE, False),
       ('   |   임계값 ', 14, GOLD, True), ('THRESHOLD = 0.5', 14, MUTE, False)]])

# ════════════════════ 9. 트러블슈팅 ════════════════════
s = slide(LIGHT)
kicker(s, Inches(0.7), Inches(0.5), '08  TROUBLESHOOTING')
text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8), [('개발 중 해결한 문제', 38, INKTXT, True)])
trbs = [
    ('@RequestParam 이름 누락', "IllegalArgumentException ('-parameters')", "어노테이션에 name 명시 → @RequestParam(name=\"min\")", RED),
    ('WTP 배포 누락', '컴파일 클래스가 wtpwebapps에 0개 배포', 'target/classes 동기화 + 서버 Clean 재배포', GOLD),
    ("SQL '&' 치환변수 오류", '리뷰 내 & 기호를 변수로 오인해 입력 멈춤', '스크립트 상단 SET DEFINE OFF 추가', BLUE),
    ('제목 링크 깨짐', "'+'·공백·괄호가 URL에서 변형 → 상세 404", '<c:url>+<c:param>으로 URL 인코딩', GREEN),
]
ty = 1.95
for name, prob, sol, col in trbs:
    card(s, Inches(0.7), Inches(ty), Inches(11.9), Inches(1.15), CARD, radius=0.10)
    circle(s, Inches(0.95), Inches(ty+0.32), Inches(0.5), col)
    text(s, Inches(1.65), Inches(ty+0.16), Inches(4.4), Inches(0.85),
         [[(name, 16, INKTXT, True)], [(prob, 12, GREY, False)]], line_spacing=1.05, space_after=2)
    text(s, Inches(6.4), Inches(ty+0.16), Inches(0.9), Inches(0.85), [('해결', 12, col, True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(7.2), Inches(ty+0.16), Inches(5.2), Inches(0.85),
         [(sol, 13, INKTXT, False)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    ty += 1.28

# ════════════════════ 9B. 한 것 / 못한 것 ════════════════════
s = slide(LIGHT)
kicker(s, Inches(0.7), Inches(0.5), 'SCOPE  ·  DONE vs TODO')
text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.8), [('구현한 기능 · 아쉬운 점(미구현)', 33, INKTXT, True)])
# 한 것
card(s, Inches(0.7), Inches(1.9), Inches(5.9), Inches(4.95), CARD, radius=0.06)
chead = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.9), Inches(5.9), Inches(0.8))
chead.fill.solid(); chead.fill.fore_color.rgb = GREEN; chead.line.fill.background(); chead.shadow.inherit = False
chead.adjustments[0] = 0.14
text(s, Inches(1.0), Inches(1.9), Inches(5.4), Inches(0.8), [('✅  구현 완료 (한 것)', 18, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
done = [
    '도서 검색 · 평점/리뷰 순위 · 도서 상세',
    '평점 분포 + 감성(긍정/중립/부정) 시각화',
    '회원가입 · 세션 로그인 · 카카오 OAuth',
    'face-api.js 얼굴 등록 · 얼굴 로그인',
    '별점 리뷰 · 좋아요 · 찜 · 마이페이지',
    '관리자 대시보드 · 회원관리 · 표지관리',
    'ON DELETE CASCADE 회원 데이터 정리',
]
text(s, Inches(1.0), Inches(2.95), Inches(5.4), Inches(3.7),
     [[(('•  '+d), 13.5, INKTXT, False)] for d in done], line_spacing=1.25, space_after=8)
# 못한 것
card(s, Inches(6.75), Inches(1.9), Inches(5.85), Inches(4.95), CARD, radius=0.06)
thead = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(1.9), Inches(5.85), Inches(0.8))
thead.fill.solid(); thead.fill.fore_color.rgb = RED; thead.line.fill.background(); thead.shadow.inherit = False
thead.adjustments[0] = 0.14
text(s, Inches(7.05), Inches(1.9), Inches(5.4), Inches(0.8), [('⚠  미구현 · 아쉬운 점 (못한 것)', 18, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
todo = [
    '관리자 기능 한정 — 회원 강제탈퇴·리뷰 신고/삭제 없음',
    '실시간 크롤링 파이프라인 미연동 (정적 데이터)',
    '감성 분석이 별점 규칙 기반 (NLP 모델 아님)',
    '검색 페이징 · 자동완성 미구현',
    '도서 표지 일부 누락 · 이미지 업로드 미지원',
    '모바일 반응형 최적화 부족',
]
text(s, Inches(7.05), Inches(2.95), Inches(5.4), Inches(3.7),
     [[(('•  '+t), 13.5, INKTXT, False)] for t in todo], line_spacing=1.25, space_after=8)

# ════════════════════ 9C. 향후 계획 (로드맵) ════════════════════
s = slide(LIGHT)
kicker(s, Inches(0.7), Inches(0.5), 'ROADMAP  ·  FUTURE WORK')
text(s, Inches(0.7), Inches(0.85), Inches(12.2), Inches(0.8), [('향후 계획 — 발전 로드맵', 33, INKTXT, True)])
text(s, Inches(0.7), Inches(1.65), Inches(12.2), Inches(0.4),
     [[('현재 아쉬운 점을 다음 단계로 이어가는 계획.  ', 13, GREY, False),
       ('★ 얼굴인식 서버 분리를 1순위', 13, BLUE, True), ('로 진행한다.', 13, GREY, False)]])
plans = [
    ('🔐  보안 · 아키텍처', BLUE, '최우선 (다음 단계)', [
        '얼굴인식 FastAPI 서버 분리',
        '서버 측 얼굴 인증 (보안 강화)',
    ]),
    ('🛠  기능 확장', GREEN, '중기', [
        '관리자 도서 CRUD (추가·삭제)',
        '도서 표지 이미지 업로드',
        '검색 페이징 · 자동완성',
    ]),
]
xs = [0.7, 6.75]
for (title_, col, tag, items), x in zip(plans, xs):
    card(s, Inches(x), Inches(2.25), Inches(5.85), Inches(4.05), CARD, radius=0.06)
    head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.25), Inches(5.85), Inches(0.95))
    head.fill.solid(); head.fill.fore_color.rgb = col; head.line.fill.background(); head.shadow.inherit = False
    head.adjustments[0] = 0.12
    text(s, Inches(x+0.35), Inches(2.4), Inches(5.2), Inches(0.45), [(title_, 18, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(x+0.35), Inches(2.82), Inches(5.2), Inches(0.35), [(tag, 12, WHITE, False)])
    text(s, Inches(x+0.4), Inches(3.55), Inches(5.1), Inches(2.5),
         [[(('•  '+it), 14.5, INKTXT, (i == 0 and col is BLUE))] for i, it in enumerate(items)],
         line_spacing=1.4, space_after=11)
text(s, Inches(0.7), Inches(6.6), Inches(12.2), Inches(0.4),
     [[('★ 교수님 피드백 반영 — ', 12.5, BLUE, True),
       ('브라우저에서 하던 얼굴인식을 독립 Python 서버(FastAPI)로 분리해 보안·확장성 개선', 12.5, GREY, False)]])

# ════════════════════ 10. 결론 (dark) ════════════════════
s = slide(INK)
kicker(s, Inches(0.7), Inches(0.6), '09  CONCLUSION')
text(s, Inches(0.7), Inches(1.0), Inches(11.5), Inches(1.0), [('결론 & 향후 과제', 40, WHITE, True)])
text(s, Inches(0.7), Inches(2.2), Inches(11.8), Inches(1.2),
     [[('Spring MVC · MyBatis · Oracle 3계층 구조로 ', 17, MUTE, False), ('도서 리뷰 분석 플랫폼', 17, GOLD, True),
       ('을 구현하고,', 17, MUTE, False)],
      [('브라우저 딥러닝(face-api.js)을 접목해 ', 17, MUTE, False), ('얼굴 인식 로그인', 17, GOLD, True),
       ('까지 완성했다.', 17, MUTE, False)]],
     line_spacing=1.3, space_after=6)
# 향후 과제 카드
todos = [('얼굴인식 서버 분리', 'FastAPI로 서버 측 인증 (보안↑)'), ('데이터 확장', '실시간 크롤링 파이프라인 연동'), ('감성 고도화', 'NLP 모델 기반 감성 분류')]
x = 0.7
for name, desc in todos:
    card(s, Inches(x), Inches(4.1), Inches(3.85), Inches(1.9), INK2, radius=0.10)
    circle(s, Inches(x+0.3), Inches(4.4), Inches(0.45), GOLD, '→', 18, INK)
    text(s, Inches(x+0.95), Inches(4.4), Inches(2.7), Inches(0.5), [(name, 17, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(x+0.35), Inches(5.15), Inches(3.2), Inches(0.7), [(desc, 13, MUTE, False)], line_spacing=1.1)
    x += 4.05
text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
     [[('BookReview', 14, GOLD, True), ('  ·  Spring 발표자료  ·  kr.ac.kopo  ·  Thank you', 13, GREY, False)]])

# ── 전체본 백업 저장 ──
prs.save(r'D:\Lecture\spring-workspace\book-review\docs\BookReview_스프링발표자료_full.pptx')

# ── 발표용: 핵심 4가지만 남기기 (표지·개요·ERD·권한별 기능·한거/못한거) ──
KEEP_ONLY_CORE = True
if KEEP_ONLY_CORE:
    keep_titles = ['ERD — 전체', '사용자 권한별', '로그인 — 3가지', '구현한 기능']  # ERD / 권한별 / 로그인 / 한거·못한거
    keep_idx = {0, 1}  # 0=표지, 1=프로젝트 개요(주제)
    for i, sl in enumerate(prs.slides):
        for sh in sl.shapes:
            if sh.has_text_frame and any(t in sh.text_frame.text for t in keep_titles):
                keep_idx.add(i); break
    sldIdLst = prs.slides._sldIdLst
    for i, sldId in enumerate(list(sldIdLst)):
        if i not in keep_idx:
            sldIdLst.remove(sldId)

prs.save(r'D:\Lecture\spring-workspace\book-review\docs\BookReview_스프링발표자료.pptx')
print('발표용 PPT 저장 — 슬라이드', len(prs.slides._sldIdLst), '(전체본: _full.pptx)')
